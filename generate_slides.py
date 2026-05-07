"""
ADInt Explorer — PPTX generator
Design guidelines applied from: Yang et al. (2025) Dashboard Vision, IEEE TVCG
  O1 — Subtitles state the key takeaway (highest saliency effort)
  O2 — Big Number widgets surface critical quantitative data
  O3 — Filter labels carry contextual meaning
  O4 — Text labels placed directly on visual marks
  L1 — Stratified top-down hierarchy: title → takeaway → stats → details
  L2 — Most important content placed upper-left (F-pattern / attention bias)
  L3 — Grouped cards use short phrases, not paragraphs

Run: pip install python-pptx && python generate_slides.py
Output: ADInt_Explorer_Presentation.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
import math

# ── Colors ────────────────────────────────────────────────────────────────────
MAROON = RGBColor(0x7A, 0x0C, 0x2E)
GOLD   = RGBColor(0xFF, 0xCC, 0x33)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
DARK   = RGBColor(0x1A, 0x1A, 0x1A)
MUTED  = RGBColor(0x55, 0x55, 0x55)
LIGHT  = RGBColor(0xF5, 0xF5, 0xF5)
BORDER = RGBColor(0xE0, 0xE0, 0xE0)

# Zhang lab color convention — must match app.py exactly
CAT_COLORS = {
    "Drug":                                 RGBColor(0xF3, 0x9C, 0x3E),
    "Dietary Supplement":                   RGBColor(0x6F, 0xBF, 0x73),
    "Complementary and Integrative Health": RGBColor(0x2C, 0xA6, 0xA4),
    "Disease or Syndrome":                  RGBColor(0xE5, 0x73, 0x6A),
    "Gene or Genome":                       RGBColor(0x4A, 0x7F, 0xD6),
    "Organism Function":                    RGBColor(0x9B, 0x59, 0xB6),
    "Other":                                RGBColor(0xB0, 0xB0, 0xB0),
}

def h(hex_str):
    return RGBColor(int(hex_str[0:2], 16), int(hex_str[2:4], 16), int(hex_str[4:6], 16))

# ── Presentation setup ────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

def slide():
    return prs.slides.add_slide(BLANK)

def box(sl, x, y, w, ht, fill, border=None, bw=0.5):
    s = sl.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE,
                            Inches(x), Inches(y), Inches(w), Inches(ht))
    s.fill.solid()
    s.fill.fore_color.rgb = fill if isinstance(fill, RGBColor) else h(fill)
    if border:
        s.line.color.rgb = border if isinstance(border, RGBColor) else h(border)
        s.line.width = Pt(bw)
    else:
        s.line.fill.background()
    return s

def oval(sl, x, y, w, ht, fill, border=None, bw=1):
    s = sl.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL,
                            Inches(x), Inches(y), Inches(w), Inches(ht))
    s.fill.solid()
    s.fill.fore_color.rgb = fill if isinstance(fill, RGBColor) else h(fill)
    if border:
        s.line.color.rgb = border if isinstance(border, RGBColor) else h(border)
        s.line.width = Pt(bw)
    else:
        s.line.fill.background()
    return s

def txt(sl, text, x, y, w, ht, size=12, bold=False, italic=False,
        color=None, align="left"):
    if color is None:
        color = DARK
    elif isinstance(color, str):
        color = h(color)
    tb = sl.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(ht))
    tf = tb.text_frame
    tf.word_wrap = True
    al = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER,
          "right": PP_ALIGN.RIGHT}.get(align, PP_ALIGN.LEFT)
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = al
        r = p.add_run()
        r.text = line
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.color.rgb = color

def lbar(sl, x, y, ht, color):
    """Thin left accent bar."""
    box(sl, x, y, 0.07, ht, color)

def hdr(sl, title, takeaway=None):
    """
    Header bar + optional O1 takeaway subtitle.
    O1: Subtitles state the slide's single key insight — they receive
    higher focused viewing effort than main titles (Yang et al., 2025).
    """
    box(sl, 0, 0, 13.33, 0.82, MAROON)
    txt(sl, title, 0.4, 0.12, 12.5, 0.65, size=22, bold=True, color=WHITE)
    if takeaway:
        box(sl, 0, 0.82, 13.33, 0.44, h("F4F1F8"))
        lbar(sl, 0, 0.82, 0.44, GOLD)
        txt(sl, takeaway, 0.3, 0.88, 12.8, 0.36, size=12, italic=True, color=MUTED)

def big_stat(sl, x, y, number, label, accent):
    """
    O2: Big Number widget — primary focal point for quantitative data.
    Numbers receive highest saliency coverage among all text types.
    """
    box(sl, x, y, 2.85, 1.45, LIGHT, accent, 1.5)
    txt(sl, number, x, y + 0.08, 2.85, 0.82, size=30, bold=True,
        color=accent, align="center")
    txt(sl, label,  x, y + 0.95, 2.85, 0.42, size=10,
        color=MUTED, align="center")

# ── Slide 1: Title ────────────────────────────────────────────────────────────
# L1: Stratified top-down — title → subtitle → big stats → context
# O2: Four key numbers anchored at bottom as primary focal points
s = slide()
box(s, 0, 0, 13.33, 7.5, h("FAFAFA"))
box(s, 0.5, 0.85, 2.2, 0.12, GOLD)
txt(s, "ADInt Explorer", 0.5, 1.05, 12, 1.3, size=46, bold=True, color=MAROON)
# O1: Subtitle as takeaway — the single thing to remember from this talk
txt(s, "Turning 1.29 million literature-backed biomedical connections into an explorable "
       "interface — no coding required.",
    0.5, 2.55, 11.5, 0.68, size=15, italic=True, color=MUTED)
txt(s, "Aviral Bhatnagar", 0.5, 3.45, 8, 0.42, size=15, bold=True)
txt(s, "Health Informatics PhD Program · University of Minnesota",
    0.5, 3.9, 10, 0.35, size=13, color=MUTED)
txt(s, "HINF 5620 — Data Visualization for the Health Sciences  |  Instructor: David Pieczkiewicz",
    0.5, 4.28, 12, 0.35, size=12, italic=True, color=MUTED)

# O2: Four big-number widgets — the numbers that define the project's scale
stats = [
    ("1,294,814", "evidence-backed edges"),
    ("162,213",   "named concepts"),
    ("7",         "intervention categories"),
    ("9",         "predicate types"),
]
for i, (n, l) in enumerate(stats):
    big_stat(s, 0.5 + i * 3.1, 5.1, n, l, MAROON)

# ── Slide 2: Problem ──────────────────────────────────────────────────────────
# L1: Stratified — most critical fact first (upper-left, L2), then context
# L3: Three short cards — one headline phrase each, no paragraph prose
s = slide()
hdr(s, "Alzheimer's has a data problem — and a visibility problem",
    takeaway="The knowledge graph exists. The problem is no one can navigate it.")

# L2: Most important fact upper-left — the 0 approved drugs
# O2: Big number as primary anchor
big_stat(s, 0.4, 1.42, "0", "FDA-approved disease-modifying drugs for AD", MAROON)
big_stat(s, 3.45, 1.42, "162,213", "concepts waiting to be explored", h("4A7FD6"))
big_stat(s, 6.5, 1.42, "1.29M", "NLP-extracted literature connections", h("6FBF73"))

# L3: Three short cards — problem → existing work → gap (stratified, top to bottom)
cards = [
    ("Drug repurposing is the fastest path",
     "Approved compounds already have safety profiles. Finding new uses is cheaper and faster.",
     "F39C3E"),
    ("Zhang lab built the knowledge graph",
     "ADInt (Xiao et al., Scientific Reports 2024) extracted 1.29M triples from biomedical text.",
     "4A7FD6"),
    ("But the KG has no interactive interface",
     "Existing output: a static file. 162K nodes in a terminal — unusable for clinicians.",
     "E5736A"),
]
for i, (title, body, c) in enumerate(cards):
    y = 3.15 + i * 1.35
    lbar(s, 0.4, y, 1.15, h(c))
    txt(s, title, 0.65, y + 0.05, 12.3, 0.38, size=13, bold=True)
    # L3: Short body — one sentence only
    txt(s, body,  0.65, y + 0.48, 12.3, 0.62, size=11, color=MUTED)

# ── Slide 3: The ADInt Knowledge Graph ───────────────────────────────────────
# L1: Top → big stats, Middle → categories + predicates, Bottom → key fact
# L2: Category legend (more visually distinctive) placed upper-left
s = slide()
hdr(s, "ADInt Knowledge Graph",
    takeaway="Every edge traces back to the exact biomedical sentence that created it.")

# O2: Two key numbers — L2 places them at top-left
big_stat(s, 0.4, 1.42, "1,294,814", "triples (PMID-level rows)", MAROON)
big_stat(s, 3.45, 1.42, "162,213", "named concepts (nodes)", h("4A7FD6"))
txt(s, "Source: Xiao, Hou, Zhou, Zhang et al. — Scientific Reports 2024",
    6.5, 1.55, 6.6, 0.35, size=11, italic=True, color=MUTED)

# L2: Node categories upper-left in middle section — most visually distinctive info
txt(s, "7 Node Categories", 0.4, 3.12, 5.5, 0.38, size=13, bold=True, color=MAROON)
for i, (name, color) in enumerate(CAT_COLORS.items()):
    oval(s, 0.44, 3.6 + i * 0.5, 0.2, 0.2, color)
    # O3: Labels carry contextual meaning — not just the name
    txt(s, name, 0.77, 3.56 + i * 0.5, 5.0, 0.35, size=11)

# Right column: predicates (secondary info)
txt(s, "9 Relationship Predicates", 6.2, 3.12, 6.8, 0.38, size=13, bold=True, color=MAROON)
rels = [
    ("TREATS",           "intervention reduces disease"),
    ("CAUSES",           "substance induces condition"),
    ("PREVENTS",         "intervention blocks outcome"),
    ("COEXISTS_WITH",    "two entities co-occur"),
    ("INTERACTS_WITH",   "compounds interact"),
    ("AUGMENTS",         "enhances another effect"),
    ("STIMULATES",       "activates a process"),
    ("ASSOCIATED_WITH",  "statistical co-occurrence"),
    ("AFFECTS",          "general influence"),
]
for i, (rel, desc) in enumerate(rels):
    col, row = i % 3, i // 3
    bx, by = 6.2 + col * 2.35, 3.6 + row * 0.68
    box(s, bx, by, 2.2, 0.58, h("EEF2FF"), h("AABDE0"), 0.5)
    txt(s, rel,  bx + 0.1, by + 0.04, 2.05, 0.28, size=9, bold=True, color="3344BB", align="center")
    # O3: Each predicate label includes what it means
    txt(s, desc, bx + 0.1, by + 0.3,  2.05, 0.22, size=7, italic=True, color=MUTED, align="center")

# O4: Key fact placed as a callout — text label near the data it describes
box(s, 0.4, 7.0, 12.53, 0.38, h("FFFBEA"), GOLD, 1)
txt(s, "Each edge stores: SUBJECT_CUI  ·  PREDICATE  ·  OBJECT_CUI  ·  PMID  ·  SENTENCE (exact quote from paper)",
    0.6, 7.05, 12.2, 0.3, size=10, color="7a5c00")

# ── Slide 4: Four Core Features ───────────────────────────────────────────────
# L1: Top-down — header → takeaway → 4 features in 2×2 grid
# L2: Evidence Panel placed upper-left — it is the key differentiator
# L3: Each card = bold label + one line + one metric (no paragraphs)
s = slide()
hdr(s, "Four Features Designed for Non-Computational Researchers",
    takeaway="Search → filter → trace → read. No code, no terminal, no installation.")

# L2: Evidence Panel upper-left — the standout feature
features = [
    # (title,                 metric,                    detail,                                      color,   col, row)
    ("Evidence Panel",
     "1,294,814 sentences",
     "Click any edge → read the exact paper sentence that created the connection.",
     "E5736A", 0, 0),
    ("Search & Explore",
     "All 162,213 concepts",
     "Type any drug, gene, or disease. The wheel re-centers around your concept.",
     "F39C3E", 1, 0),
    ("Mechanistic Path Finder",
     "Shortest path in <1 s",
     "Pick any two concepts — see the chain of predicates that connects them.",
     "4A7FD6", 0, 1),
    ("Category Filter",
     "7 toggle-able types",
     "Hide genes, show only drugs and dietary supplements — one checkbox each.",
     "6FBF73", 1, 1),
]
for title, metric, detail, c, col, row in features:
    x = 0.4 + col * 6.5
    y = 1.45 + row * 2.82
    box(s, x, y, 6.15, 2.65, LIGHT, h(c), 1.5)
    lbar(s, x, y, 2.65, h(c))
    # O2: Metric as big number inside the card
    txt(s, metric, x + 0.25, y + 0.12, 5.65, 0.48, size=18, bold=True, color=h(c))
    txt(s, title,  x + 0.25, y + 0.65, 5.65, 0.4,  size=13, bold=True)
    # L3: One sentence only — no paragraphs
    txt(s, detail, x + 0.25, y + 1.1,  5.65, 0.9,  size=11, color=MUTED)

# ── Slide 5: Evidence Panel ───────────────────────────────────────────────────
# L1: Top-down — takeaway → big number → mock-up → features
# L2: Big number upper-left; mock-up panel takes visual center
# O4: Mock-up labels placed directly adjacent to the elements they describe
s = slide()
hdr(s, "Evidence Panel — The Key Differentiator",
    takeaway="No existing ADInt visualization shows the sentence behind the connection. This one does.")

# O2: Big number upper-left (L2)
big_stat(s, 0.4, 1.42, "1,294,814", "literature sentences surfaced on demand", MAROON)

# O1: The one fact the audience must remember
txt(s, "Click any edge. Read the paper.", 3.55, 1.58, 9.4, 0.55, size=20, bold=True, color=MAROON)
txt(s, "No other ADInt tool offers sentence-level evidence — only CUI identifiers.",
    3.55, 2.18, 9.4, 0.45, size=12, italic=True, color=MUTED)

# Mock-up of the evidence panel (left side)
box(s, 0.4, 3.1, 5.6, 4.0, h("FFFFFF"), BORDER, 1)
txt(s, "Evidence Panel", 0.55, 3.2, 5.3, 0.38, size=13, bold=True)

txt(s, "TREATS", 0.55, 3.72, 5.0, 0.3, size=11, italic=True, color=MUTED)
lbar(s, 0.58, 4.1, 1.05, BORDER)
txt(s, "PMID 28945394", 0.78, 4.12, 4.8, 0.3, size=11, color="4A7FD6")
txt(s, '"Curcumin reduces amyloid plaque formation in Alzheimer\'s '
       'disease animal models via neuroprotective mechanisms."',
    0.78, 4.46, 4.85, 0.7, size=10, italic=True, color=MUTED)

txt(s, "COEXISTS_WITH", 0.55, 5.3, 5.0, 0.3, size=11, italic=True, color=MUTED)
lbar(s, 0.58, 5.68, 1.05, BORDER)
txt(s, "PMID 1394115", 0.78, 5.7, 4.8, 0.3, size=11, color="4A7FD6")
txt(s, '"Turmeric and curcumin reversed aflatoxin-induced liver '
       'damage in a controlled feeding study."',
    0.78, 6.04, 4.85, 0.62, size=10, italic=True, color=MUTED)

# O4: Labels placed directly beside the elements they describe (not in a separate legend)
txt(s, "← predicate type", 5.75, 3.72, 2.8, 0.3, size=9, italic=True, color="9B59B6")
txt(s, "← clickable PubMed link", 5.75, 4.12, 2.8, 0.3, size=9, italic=True, color="9B59B6")
txt(s, "← exact paper sentence", 5.75, 4.7,  2.8, 0.3, size=9, italic=True, color="9B59B6")

# Right column: three key properties (L3: short phrases)
ev_props = [
    ("Two trigger modes",   "Edge click → pair evidence.  Path finder → every hop in chain.",      "E5736A"),
    ("Live PubMed links",   "Every PMID links to pubmed.ncbi.nlm.nih.gov — opens in browser.",     "4A7FD6"),
    ("Research instrument", "Transforms the browser from a visualization into a literature tool.", "6FBF73"),
]
for i, (t, b, c) in enumerate(ev_props):
    y = 3.1 + i * 1.32
    lbar(s, 6.5, y, 1.1, h(c))
    txt(s, t, 6.72, y + 0.06, 6.5, 0.38, size=13, bold=True)
    txt(s, b, 6.72, y + 0.48, 6.5, 0.6,  size=11, color=MUTED)

# ── Slide 6: Mechanistic Path Finder ─────────────────────────────────────────
# L1: Top-down — takeaway → diagram (with O4 inline labels) → stats
# O4: Path node labels sit directly ON the ovals, not in a separate callout
s = slide()
hdr(s, "Mechanistic Path Finder — Trace How Treatments Reach Alzheimer's Disease",
    takeaway="Shortest path across 742,274 graph edges — computed in under 1 second.")

# O4: Path diagram — labels directly ON the nodes (not in a separate legend box)
# Node positions: left-to-right linear path, vertically centered
path_items = [
    ("Photodynamic\nTherapy",   "2CA6A4",  0.55, 3.1),
    ("Immune\nResponse",        "9B59B6",  3.55, 3.1),
    ("Alzheimer's\nDisease",    "E5736A",  6.55, 3.1),
]
for lbl, c, px, py in path_items:
    oval(s, px, py, 1.85, 1.25, h(c), GOLD, 2)
    # O4: Label is ON the node, not beside it
    txt(s, lbl, px, py + 0.28, 1.85, 0.68, size=10, bold=True, color=WHITE, align="center")

# Predicate labels positioned ON the arrows (O4)
txt(s, "STIMULATES", 2.48, 3.32, 1.15, 0.4, size=9, bold=True, color="7a5c00", align="center")
box(s, 2.43, 3.52, 1.1, 0.14, h("FFF0A0"))  # arrow line
txt(s, "►",          3.32, 3.43, 0.3, 0.28, size=10, bold=True, color="7a5c00")

txt(s, "ASSOCIATED\nWITH", 5.47, 3.2, 1.16, 0.52, size=9, bold=True, color="7a5c00", align="center")
box(s, 5.42, 3.52, 1.1, 0.14, h("FFF0A0"))
txt(s, "►",               6.33, 3.43, 0.3, 0.28, size=10, bold=True, color="7a5c00")

# O2: Big stat placed upper-left (L2) — the claim the audience should remember
big_stat(s, 0.4, 1.42, "742,274", "undirected graph edges searched", h("4A7FD6"))
big_stat(s, 3.45, 1.42, "< 1 s", "NetworkX shortest-path computation", h("9B59B6"))

# Path string callout (L1 — placed below the diagram)
box(s, 0.4, 4.62, 8.8, 0.52, h("FFFBEA"), GOLD, 1)
txt(s, "Path: Photodynamic Therapy → Immune Response → Alzheimer's Disease  (2 hops)",
    0.6, 4.69, 8.4, 0.38, size=10, bold=True, color="7a5c00")

# Right column: what path finding enables (L3: short phrases)
path_props = [
    ("Gold highlighting",   "Path nodes and edges render in gold — survive all category filters.",  "F39C3E"),
    ("Auto-populates panel","Path finder fills evidence panel for every hop — no extra clicks.",    "E5736A"),
    ("Status bar readout",  "Full path string shown in status bar after each search.",              "4A7FD6"),
]
for i, (t, b, c) in enumerate(path_props):
    y = 1.42 + i * 1.35
    lbar(s, 9.45, y, 1.12, h(c))
    txt(s, t, 9.68, y + 0.06, 3.5, 0.38, size=12, bold=True)
    txt(s, b, 9.68, y + 0.48, 3.5, 0.62, size=10, color=MUTED)

# ── Slide 7: Design Rationale ─────────────────────────────────────────────────
# Meta-slide: connects our choices to Dashboard Vision guidelines
# L1: Three-row stratified layout — one guideline per row
# O2: Guideline codes (O1, O2, L1…) act as number anchors
s = slide()
hdr(s, "Design Choices Grounded in Visualization Science",
    takeaway="Each decision maps to an empirical attention finding from Yang et al. (2025) IEEE TVCG.")

txt(s, "Dashboard Vision Guidelines Applied in ADInt Explorer",
    0.4, 1.35, 12.5, 0.42, size=14, bold=True, color=MAROON)

rows = [
    ("L1", "Stratified layout maximizes\nsaliency coverage",
     "Category-wedge wheel",
     "7 equal angular wedges guarantee all node types appear — rare categories\n"
     "(CIH, Organism Function) are never visually crushed by high-degree nodes.",
     "4A7FD6"),
    ("O2", "Big Numbers are the\nprimary focal point",
     "Status bar + stat widgets",
     "Edge count, path length, and per-category counts displayed as large\n"
     "text anchors — the first thing the eye lands on after the focal node.",
     "F39C3E"),
    ("O1 + O4", "Subtitles & inline labels\ndrive the deepest attention",
     "Evidence panel text",
     "The exact source sentence and PMID placed directly beside the edge —\n"
     "not in a tooltip or legend. Text close to a mark gets the highest SC.",
     "E5736A"),
    ("L3", "Simplified grouped views\nkeep focus on main content",
     "Focal-only edge toggle",
     "Default ON: hides cross-links between neighbor nodes. 253 edges → 69.\n"
     "Reduces clutter 73% while preserving all paths and focal connections.",
     "6FBF73"),
]
for i, (code, principle, our_choice, detail, c) in enumerate(rows):
    y = 1.92 + i * 1.34
    # O2: Guideline code as large anchor (upper-left within each row, L2)
    box(s, 0.4, y, 1.0, 1.18, h(c), h(c), 1)
    txt(s, code, 0.4, y + 0.28, 1.0, 0.62, size=16, bold=True, color=WHITE, align="center")
    # Principle (what the paper found)
    txt(s, principle, 1.55, y + 0.06, 3.0, 0.62, size=11, bold=True, color=h(c))
    # Our choice (what we built)
    box(s, 4.72, y + 0.08, 2.2, 0.48, h("F0F0F0"), BORDER, 0.5)
    txt(s, our_choice, 4.82, y + 0.14, 2.0, 0.35, size=10, bold=True, align="center")
    # L3: Short detail — two lines max
    txt(s, detail, 7.1, y + 0.06, 5.95, 0.9, size=10, color=MUTED)

txt(s, "Yang, Hou, Li, Chang, & Zeng (2025). Dashboard Vision. IEEE TVCG Vol. 31 No. 10.",
    0.4, 7.08, 12.5, 0.3, size=9, italic=True, color=MUTED)

# ── Slide 8: Conclusions & Future Work ───────────────────────────────────────
# L1: Top-down — tagline → contributions (left) → future work (right)
# L2: Contributions upper-left (what exists), future work upper-right (what's next)
# O2: Key stats as anchors; L3: short bullet phrases
s = slide()
box(s, 0, 0, 13.33, 7.5, h("FAFAFA"))
box(s, 0.5, 0.7, 2.5, 0.12, GOLD)
txt(s, "ADInt had the data.\nNow it has an interface.",
    0.5, 0.9, 12, 1.65, size=36, bold=True, color=MAROON)

# L2: Contributions on LEFT (what was built — primary) — upper-left anchor
txt(s, "What was built", 0.5, 2.72, 5.8, 0.4, size=13, bold=True, color=MAROON)
contributions = [
    ("Search & refocus",  "Any of 162,213 concepts as focal node",        "F39C3E"),
    ("Category filters",  "7 toggle-able node types, filter-proof anchor", "6FBF73"),
    ("Mechanistic paths", "NetworkX shortest path with gold highlighting", "4A7FD6"),
    ("Evidence panel",    "PMID + source sentence on every edge click",    "E5736A"),
]
for i, (t, b, c) in enumerate(contributions):
    y = 3.18 + i * 0.92
    lbar(s, 0.5, y, 0.78, h(c))
    txt(s, t, 0.72, y + 0.04, 2.5, 0.35, size=12, bold=True)
    # L3: One line only
    txt(s, b, 0.72, y + 0.42, 5.5, 0.32, size=10, color=MUTED)

# Right: Future work
txt(s, "What comes next", 6.8, 2.72, 6.3, 0.4, size=13, bold=True, color=MAROON)
futures = [
    ("01", "DrKGC integration",
     "Overlay link-prediction confidence scores as edge weights",  "4A7FD6"),
    ("02", "ClinicalTrials.gov",
     "Flag nodes in active Alzheimer's trials via CT.gov API",     "6FBF73"),
    ("03", "Subgraph export",
     "Download visible subgraph as CSV/JSON for external tools",   "F39C3E"),
    ("04", "Public deployment",
     "Host on Render/Railway — browser access, no Python needed",  "9B59B6"),
]
for i, (n, t, b, c) in enumerate(futures):
    y = 3.18 + i * 0.92
    oval(s, 6.8, y + 0.06, 0.58, 0.58, h(c), h(c), 1)
    txt(s, n, 6.8, y + 0.15, 0.58, 0.38, size=10, bold=True, color=WHITE, align="center")
    txt(s, t, 7.55, y + 0.04, 5.5, 0.35, size=12, bold=True)
    # L3: One line only
    txt(s, b, 7.55, y + 0.42, 5.5, 0.32, size=10, color=MUTED)

# Closing bar
box(s, 2.5, 7.05, 8.33, 0.72, MAROON)
txt(s, "Thank you — happy to demo live or take questions",
    2.5, 7.16, 8.33, 0.5, size=14, bold=True, color=WHITE, align="center")

# ── Save ──────────────────────────────────────────────────────────────────────
prs.save("ADInt_Explorer_Presentation.pptx")
print("Saved: ADInt_Explorer_Presentation.pptx")
print(f"  8 slides total")
print("  Guidelines applied: O1 O2 O3 O4 L1 L2 L3 (Yang et al. 2025)")
